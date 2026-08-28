"""PPTX ingestion into a persistent, incremental knowledge base.

The user stores past course decks as .pptx in inputs/past_ppts/. We extract
each deck ONCE into knowledge_base/decks/<key>.json and record a hash in
knowledge_base/manifest.json. On every later run we only (re)process decks
whose file hash changed or that are new — already-ingested decks are kept as
they are and never re-extracted. Nothing from the past is dropped.

Each deck record holds:
  - a structural summary (deck title + per-slide titles)  -> always injected
  - full per-slide text + notes + tables (chunks)         -> for RAG retrieval

No API is needed for ingestion; it is pure text extraction, so the memory is
built and persisted regardless of whether the generation key is set.
"""
from __future__ import annotations
import glob
import hashlib
import json
import math
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from . import config

KB_DIR = config.KB_DIR
DECKS_DIR = KB_DIR / "decks"
# The pre-course-scoping global manifest. Kept ONLY so migrate_legacy_decks can read the
# hashes out of it; nothing writes it any more.
MANIFEST = KB_DIR / "manifest.json"


# --------------------------------------------------------------------------- #
# THE STORE IS SCOPED BY COURSE.
#
# Decks used to live at knowledge_base/decks/session_07.json — keyed by session NUMBER
# alone, one directory, one manifest, globbed globally. Invisible on a single-course
# instance, and wrong the moment there are two:
#
#   · two courses that both have a session 7 shared ONE file, so ingesting the second
#     silently overwrote the first, and both courses then read whichever was fetched
#     last as "what I have already taught";
#   · taught_digest() — the "do not teach this again" block in every generation prompt —
#     was built from whatever was on disk, so a React doc could be told it had already
#     covered Deadlock Detection;
#   · taught_titles() feeds the DETERMINISTIC repetition guardrail, so a legitimate
#     slide in one course could be failed for repeating another course's title;
#   · deleting or renumbering a session in one course moved another course's decks.
#
# Every function below therefore takes the COURSE FIRST and requires it. Deliberately
# not defaulted to the instance-wide active course: that global is exactly what made
# this wrong, and a missing argument should be a loud TypeError rather than a silent
# read of somebody else's decks.
# --------------------------------------------------------------------------- #
def course_slug(course: str) -> str:
    """A stable, filesystem-safe folder name for a course.

    Readable prefix plus a short digest of the FULL name. The digest is what makes it
    safe: 'C++ / Advanced' and 'C   Advanced' both reduce to the same readable slug, and
    two courses sharing a folder is the bug this module is being changed to fix.
    """
    name = (course or "").strip() or "default"
    base = re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_")[:48] or "course"
    return f"{base}_{hashlib.md5(name.encode('utf-8')).hexdigest()[:6]}"


def course_decks_dir(course: str) -> Path:
    return DECKS_DIR / course_slug(course)


def prereq_decks_dir(course: str, prereq: str) -> Path:
    """Where an EXTERNAL prerequisite's decks live: inside the course that declared it.

    An external prerequisite is a name and a set of slides — a course taught somewhere
    else, with no course of its own in this agent to hang decks on. So they belong to the
    course that declared it, and go when it goes.

    A subfolder rather than a sibling, because the two must never be confused: the decks
    in `decks/<course>/` are what this course has ALREADY TAUGHT (repeating one is a
    failure), and the decks in `decks/<course>/prereq/<name>/` are what the learner knew
    before session 1 (referring to one is correct). Same shape, opposite rule.
    """
    return course_decks_dir(course) / "prereq" / course_slug(prereq)


def _store_dir(course: str, prereq: str | None = None) -> Path:
    return prereq_decks_dir(course, prereq) if prereq else course_decks_dir(course)


def deck_path(course: str, session_no: int, prereq: str | None = None) -> Path:
    return _store_dir(course, prereq) / f"session_{int(session_no):02d}.json"


def _manifest_path(course: str, prereq: str | None = None) -> Path:
    return _store_dir(course, prereq) / "manifest.json"


def kb_rel(course: str, session_no: int) -> str:
    """The deck's KB-relative path, which is the key the cloud mirror stores it under."""
    return f"decks/{course_slug(course)}/session_{int(session_no):02d}.json"


# --------------------------------------------------------------------------- #
# helpers
# --------------------------------------------------------------------------- #
def _file_hash(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _session_no(path: Path) -> int | None:
    m = re.search(r"(\d+)", path.stem)
    return int(m.group(1)) if m else None


def _deck_key(path: Path) -> str:
    n = _session_no(path)
    return f"session_{n:02d}" if n is not None else re.sub(r"\W+", "_", path.stem.lower())


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())


def _slide_notes(slide) -> str:
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


def _slide_tables(shape) -> list[list[list[str]]]:
    tables = []
    if shape.has_table:
        t = shape.table
        rows = [[cell.text.strip() for cell in row.cells] for row in t.rows]
        tables.append(rows)
    return tables


# --------------------------------------------------------------------------- #
# extraction
# --------------------------------------------------------------------------- #
def extract_deck(path: Path) -> dict:
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts, tables = [], []
        for shape in slide.shapes:
            if shape.has_table:
                tables += _slide_tables(shape)
                continue
            txt = _shape_text(shape)
            if not txt:
                continue
            # first placeholder-ish text becomes the title
            is_title = getattr(shape, "is_placeholder", False) and \
                getattr(shape.placeholder_format, "idx", None) == 0
            if is_title and not title:
                title = txt.split("\n")[0].strip()
            else:
                body_parts.append(txt)
        if not title and body_parts:
            title = body_parts[0].split("\n")[0][:80]
        slides.append({
            "n": i,
            "title": title,
            "body": "\n".join(body_parts).strip(),
            "notes": _slide_notes(slide),
            "tables": tables,
        })

    deck_title = slides[0]["title"] if slides else path.stem
    summary_lines = [f"    - Slide {s['n']}: {s['title']}" for s in slides if s["title"]]
    summary = f"{deck_title}\n" + "\n".join(summary_lines)

    return {
        "session_no": _session_no(path),
        "source_file": path.name,
        "deck_title": deck_title,
        "n_slides": len(slides),
        "summary": summary,
        "slides": slides,
    }


# --------------------------------------------------------------------------- #
# persistent KB
# --------------------------------------------------------------------------- #
def _load_manifest(course: str, prereq: str | None = None) -> dict:
    path = _manifest_path(course, prereq)
    if path.exists():
        try:
            return json.loads(path.read_text())
        except Exception:
            return {}
    return {}


def _save_manifest(course: str, m: dict, prereq: str | None = None):
    d = _store_dir(course, prereq)
    d.mkdir(parents=True, exist_ok=True)
    (d / "manifest.json").write_text(json.dumps(m, indent=2), encoding="utf-8")


def put_deck(course: str, session_no: int, deck: dict,
             prereq: str | None = None) -> Path:
    """Write one extracted deck into this course's store, manifest included.

    The single place a deck is written. sync used to build the path itself, which is how
    the course scoping went missing there.
    """
    d = _store_dir(course, prereq)
    d.mkdir(parents=True, exist_ok=True)
    path = deck_path(course, session_no, prereq)
    path.write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
    manifest = _load_manifest(course, prereq)
    manifest[f"session_{int(session_no):02d}"] = {
        "hash": deck.get("source_hash") or "",
        "source_file": deck.get("source_file") or path.name,
        "session_no": deck.get("session_no", int(session_no)),
        "n_slides": deck.get("n_slides"),
    }
    _save_manifest(course, manifest, prereq)
    # MIRRORED NOW, not at the end of the run. On an ephemeral host the file just
    # written sits on a disk that goes with the next spin-down or redeploy, and the only
    # thing that copied it to the cloud DB was kb_backup() — called once, after every
    # link in a sync or a prerequisite import had been read. So a 29-link read that lost
    # its instance at link 9 lost all nine decks, while the prerequisite row (committed
    # up front) survived: a prerequisite attached to nothing. One ~50 KB write per deck
    # ends that. Best effort — a storage hiccup must not fail an extraction that worked.
    try:
        from . import db
        db.kb_put_rel(_kb_rel_path(path))
        db.kb_put_rel(_kb_rel_path(_manifest_path(course, prereq)))
    except Exception:
        pass
    return path


def _kb_rel_path(path: Path) -> str:
    """A KB-relative posix path — the key the cloud mirror stores a file under."""
    try:
        return path.resolve().relative_to(config.KB_DIR.resolve()).as_posix()
    except Exception:
        return path.name


def get_deck(course: str, session_no: int, prereq: str | None = None) -> dict | None:
    path = deck_path(course, session_no, prereq)
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text())
    except Exception:
        return None


def has_deck(course: str, session_no: int, prereq: str | None = None) -> bool:
    return deck_path(course, session_no, prereq).exists()


@dataclass
class IngestReport:
    ingested: list[str]
    skipped: list[str]
    total_decks: int


def ingest(course: str, verbose: bool = True) -> IngestReport:
    """Incrementally sync inputs/past_ppts/ into ONE course's knowledge base.

    The offline path (local .pptx files), used by the eval harness and as the fallback
    when nothing has been synced. It takes the course for the same reason everything
    else here does: the decks it writes are that course's memory, not the instance's.
    """
    d = course_decks_dir(course)
    d.mkdir(parents=True, exist_ok=True)
    pattern = config.harness()["context"]["past_ppts_glob"]
    paths = sorted(Path(p) for p in glob.glob(str(config.ROOT / pattern)))
    manifest = _load_manifest(course)

    ingested, skipped = [], []
    for path in paths:
        key = _deck_key(path)
        fhash = _file_hash(path)
        rec = manifest.get(key)
        deck_json = d / f"{key}.json"
        if rec and rec.get("hash") == fhash and deck_json.exists():
            skipped.append(path.name)   # already in memory, unchanged
            continue
        deck = extract_deck(path)
        deck_json.write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
        manifest[key] = {
            "hash": fhash,
            "source_file": path.name,
            "session_no": deck["session_no"],
            "n_slides": deck["n_slides"],
        }
        ingested.append(path.name)

    _save_manifest(course, manifest)
    if verbose:
        print(f"[KB] {course}: ingested {len(ingested)} new/changed deck(s), "
              f"skipped {len(skipped)} cached, {len(manifest)} total in memory.")
    return IngestReport(ingested, skipped, len(manifest))


def load_all_decks(course: str, prereq: str | None = None) -> list[dict]:
    decks = []
    # A non-recursive glob, deliberately: decks/<course>/prereq/… is a DIFFERENT store
    # with the opposite rule attached to it, and sweeping it in here would make a
    # prerequisite's topics look like this course's own prior sessions.
    for f in sorted(_store_dir(course, prereq).glob("*.json")):
        if f.name == "manifest.json":
            continue
        try:
            decks.append(json.loads(f.read_text()))
        except Exception:
            continue
    decks.sort(key=lambda d: (d.get("session_no") is None, d.get("session_no") or 0))
    return decks


def deck_session_numbers(course: str, prereq: str | None = None) -> set[int]:
    """Which of this course's sessions have an extracted deck — from the FILENAMES,
    nothing parsed.

    Callers that only need "does this session have a deck?" were using load_all_decks(),
    which reads and JSON-parses every deck in the course (1.1 MB across 30 decks here)
    and was being called three times on a single page load. The name carries the answer.
    """
    out = set()
    for p in _store_dir(course, prereq).glob("session_*.json"):
        m = re.search(r"session_(\d+)\.json$", p.name)
        if m:
            out.add(int(m.group(1)))
    return out


def courses_with_decks() -> list[str]:
    """The course SLUGS that have a deck folder. For the migration and for diagnostics —
    a slug cannot be turned back into a course name, so callers that need names match
    these against course_slug() of the courses they know about."""
    if not DECKS_DIR.is_dir():
        return []
    return sorted(d.name for d in DECKS_DIR.iterdir()
                  if d.is_dir() and any(d.glob("session_*.json")))


UNASSIGNED = "_unassigned"


def legacy_decks() -> dict:
    """{session_no: path} for decks still in the OLD flat layout, if any.

    `decks/session_NN.json` — written before the store was scoped by course. Present on
    every instance that ran an earlier version.
    """
    out = {}
    if not DECKS_DIR.is_dir():
        return out
    for f in DECKS_DIR.glob("session_*.json"):
        if not f.is_file():
            continue
        m = re.search(r"session_(\d+)\.json$", f.name)
        if m:
            out[int(m.group(1))] = f
    return out


def migrate_legacy_decks() -> dict:
    """Move flat `decks/session_NN.json` files into the folder of the course that owns
    them. Runs once at startup; a no-op afterwards.

    WHICH COURSE OWNS A DECK is inferred, because the old layout never recorded it. In
    order of confidence:
      1. exactly one course whose curriculum has that session number WITH a deck link —
         a deck exists on disk because a link was extracted, so this is the strong signal;
      2. exactly one course that merely has that session number;
      3. otherwise AMBIGUOUS: parked under decks/_unassigned/ and named in the return
         value. Deliberately not guessed — attributing one course's material to another
         is worse than leaving it aside, because the writer would then be told it had
         already taught something it had not.

    The per-course manifests are rebuilt from the legacy global manifest so an already
    extracted deck is still recognised as extracted and is not re-downloaded.

    Returns {"moved": {session_no: course}, "unassigned": [session_no], "kb_paths": [...]}
    — `kb_paths` being the new KB-relative paths, so the caller can repoint the cloud
    mirror in one statement instead of one per deck.
    """
    legacy = legacy_decks()
    if not legacy:
        return {"moved": {}, "unassigned": [], "kb_paths": []}

    from . import db
    # ONE read of the curriculum for the whole migration rather than one per deck.
    rows = []
    try:
        rows = db._query("SELECT course, session_no, ppt_link FROM curriculum")
    except Exception as e:
        print(f"[decks] migration could not read the curriculum ({e!r}) — "
              f"leaving the old layout alone.")
        return {"moved": {}, "unassigned": sorted(legacy), "kb_paths": []}

    linked: dict = {}        # session_no -> {courses that have it WITH a link}
    present: dict = {}       # session_no -> {courses that have it at all}
    for r in rows:
        c, n = r.get("course"), r.get("session_no")
        if not c or n is None:
            continue
        present.setdefault(int(n), set()).add(c)
        if (r.get("ppt_link") or "").strip():
            linked.setdefault(int(n), set()).add(c)

    legacy_manifest = {}
    if MANIFEST.exists():
        try:
            legacy_manifest = json.loads(MANIFEST.read_text())
        except Exception:
            legacy_manifest = {}

    moved, unassigned, kb_paths = {}, [], []
    touched_manifests: dict = {}
    for no, src in sorted(legacy.items()):
        owners = linked.get(no) or set()
        if len(owners) != 1:
            owners = present.get(no) or set()
        if len(owners) == 1:
            course = next(iter(owners))
            dest_dir = course_decks_dir(course)
            rel = kb_rel(course, no)
        else:
            course = None
            dest_dir = DECKS_DIR / UNASSIGNED
            rel = f"decks/{UNASSIGNED}/session_{no:02d}.json"
        dest_dir.mkdir(parents=True, exist_ok=True)
        dest = dest_dir / f"session_{no:02d}.json"
        try:
            src.replace(dest)          # atomic within one filesystem
        except Exception as e:
            print(f"[decks] could not move session {no} ({e!r}) — left in place.")
            continue
        kb_paths.append(rel)
        if course:
            moved[no] = course
            key = f"session_{no:02d}"
            rec = legacy_manifest.get(key)
            if rec:
                touched_manifests.setdefault(course, _load_manifest(course))[key] = rec
        else:
            unassigned.append(no)

    for course, manifest in touched_manifests.items():
        _save_manifest(course, manifest)
    if MANIFEST.exists() and not legacy_decks():
        # The global manifest described the flat layout and nothing reads it now.
        try:
            MANIFEST.replace(KB_DIR / "manifest.legacy.json")
        except Exception:
            pass
    return {"moved": moved, "unassigned": sorted(unassigned), "kb_paths": kb_paths}


def drop_deck(course: str, session_no: int) -> bool:
    """Delete one session's extracted deck, in ONE course. True if there was one.

    Used when a session is REMOVED from the curriculum. sync.prune_orphan_decks cannot
    do this: it deliberately only touches sessions the curriculum still lists, so that a
    deck belonging to an unlisted session is never assumed to be rubbish. Here the row
    has just been deleted on purpose, so the deck is unambiguously orphaned — and if it
    were left behind, the next session to take that number would inherit it as material
    it had "already taught".
    """
    path = deck_path(course, session_no)
    existed = path.exists()
    path.unlink(missing_ok=True)
    manifest = _load_manifest(course)
    if manifest.pop(f"session_{int(session_no):02d}", None) is not None:
        _save_manifest(course, manifest)
    return existed


def drop_course_decks(course: str) -> list[int]:
    """Delete EVERY deck this course holds, and the folder itself. Returns the sessions.

    What deleting a course means now that the store is scoped. It used to be a per-session
    calculation — "drop the decks whose session number no other course still claims" —
    which existed only because one directory held every course's decks. That was wrong in
    both directions: a deck belonging to the deleted course was KEPT whenever some other
    course happened to have the same session number, and the number was all there was to
    go on. A course's decks are its own, in its own folder, and they go with it.
    """
    d = course_decks_dir(course)
    if not d.is_dir():
        return []
    gone = sorted(deck_session_numbers(course))
    for f in d.glob("*.json"):
        f.unlink(missing_ok=True)
    # …and any EXTERNAL prerequisite decks this course declared: they belong to it,
    # because there is no course of their own to hang them on.
    import shutil
    shutil.rmtree(d / "prereq", ignore_errors=True)
    try:
        d.rmdir()
    except OSError:
        pass          # something else is in there; the decks are gone, which is the point
    return gone


def drop_prereq_decks(course: str, prereq: str) -> list[int]:
    """Delete an external prerequisite's decks. Nothing else owns them."""
    d = prereq_decks_dir(course, prereq)
    if not d.is_dir():
        return []
    gone = sorted(deck_session_numbers(course, prereq))
    for f in d.glob("*.json"):
        f.unlink(missing_ok=True)
    try:
        d.rmdir()
        d.parent.rmdir()          # the `prereq/` holder, once it is empty
    except OSError:
        pass
    return gone


def renumber_decks(course: str, mapping: dict) -> list[str]:
    """Move ONE course's extracted decks to follow their sessions. Returns what moved.

    A deck lives at knowledge_base/decks/<course>/session_NN.json, keyed within the
    course by session number — so renumbering the curriculum without moving these would leave
    Session 6 reading Session 5's deck as "what I already taught". `mapping` is
    {old_session_no: new_session_no}.

    Written in TWO PASSES through temporary names. A shift like {5:6, 6:7} applied in
    place would have 5 overwrite 6 before 6 had moved, destroying a deck the user paid
    to download; parking every affected deck under a temp name first makes the order of
    operations irrelevant.
    """
    if not mapping:
        return []
    d = course_decks_dir(course)
    manifest = _load_manifest(course)
    moved: list[str] = []
    staged: list[tuple[Path, int]] = []      # (temp path, new session number)

    for old, new in mapping.items():
        if int(old) == int(new):
            continue
        src = deck_path(course, old)
        if not src.exists():
            continue
        tmp = d / f".renumber_{int(old):02d}_to_{int(new):02d}.json"
        src.replace(tmp)
        staged.append((tmp, int(new)))
        manifest.pop(f"session_{int(old):02d}", None)

    for tmp, new in staged:
        dest = deck_path(course, new)
        try:
            deck = json.loads(tmp.read_text())
            deck["session_no"] = new          # the number is inside the file too
            dest.write_text(json.dumps(deck, ensure_ascii=False, indent=2),
                            encoding="utf-8")
            tmp.unlink()
            manifest[f"session_{new:02d}"] = {
                "hash": deck.get("source_hash") or "",
                "source_file": deck.get("source_file") or dest.name,
                "session_no": new,
                "n_slides": deck.get("n_slides"),
            }
            moved.append(f"session_{new:02d}")
        except Exception:
            # Leave the temp file rather than losing the deck; a re-fetch can replace it.
            continue
    _save_manifest(course, manifest)
    return moved


def decks_before(course: str, session_no: int, prereq: str | None = None) -> list[dict]:
    return [d for d in load_all_decks(course, prereq)
            if d.get("session_no") is not None and d["session_no"] < session_no]


# --------------------------------------------------------------------------- #
# WHAT HAS ALREADY BEEN TAUGHT  (the point of ingesting the decks at all)
#
# The generator used to receive each prior deck's raw `summary` — deck title plus
# EVERY slide title, in order. Measured on this course that is 38,000 characters in
# which "Data Representation" appears eight times and "Recap"/"Quiz Time!" dozens of
# times: the model was paying for ~10k tokens of noise, and what it actually needed —
# a clean list of the TOPICS already covered, per session — was buried in it.
#
# taught_index() is that list: per prior session, its slide titles de-duplicated and
# stripped of deck furniture. 950 titled slides across this course collapse to 281
# distinct topics, so the block gets smaller AND says something. Slide BODIES are not
# summarised here — body-level detail arrives through retrieve(), which is targeted at
# the topic being written rather than dumped wholesale.
# --------------------------------------------------------------------------- #
# Deck furniture: structural slides that say nothing about what was taught.
# The second group is FURNITURE OF A DIFFERENT KIND: a slide titled bare "Overview" or
# "Examples" is a structural heading under whatever came before it, not a topic the
# learner was taught. Left in, they became "already taught" entries — so a prerequisite
# index reported that the learner knows "Overview" and "Analogy", and any takeaway using
# either word looked like a repeat. Anchored to the WHOLE title, so "Overview of Paging"
# and "Examples of Deadlock" are still topics and still counted.
_BOILERPLATE = re.compile(
    r"^(agenda|agenda for today.?s session|recap|quiz ?time!?|quiz|thank ?you|"
    r"key ?takeaways?|takeaways?|summary|questions?\??|q ?& ?a|poll|break|"
    r"upcoming session|next session|welcome|introduction|"
    r"overviews?|examples?|analogy|analogies|comparisons?|differences?|"
    r"advantages|disadvantages|pros ?(?:& ?|and )?cons|benefits|limitations|"
    r"applications|use ?cases?|case ?stud(?:y|ies)|conclusions?|objectives?|"
    r"outline|contents|references?|demo|practice|exercises?|assignments?|"
    r"problem ?statement|aspects?|so far|whats? next|lets? begin|"
    r"hands.?on|discussion|activity)\W*$", re.I)


def _clean_title(t: str) -> str:
    # Titles carry vertical tabs and newlines from PowerPoint text boxes.
    return re.sub(r"\s+", " ", str(t or "").replace("\x0b", " ")).strip(" -–—:")


def taught_index(course: str, before_session: int,
                 prereq: str | None = None) -> list[dict]:
    """Per prior session OF THIS COURSE: the distinct topics its deck actually taught.

    Returns [{session_no, deck_title, topics: [...]}] in session order, oldest first.
    """
    out = []
    for deck in decks_before(course, before_session, prereq):
        deck_title = _clean_title(deck.get("deck_title"))
        seen, topics = set(), []
        for s in deck.get("slides") or []:
            t = _clean_title(s.get("title"))
            if not t or _BOILERPLATE.match(t):
                continue
            key = t.lower()
            # The deck title repeats on section-divider slides; it is already the
            # heading of this block, so it adds nothing as a topic.
            if key == deck_title.lower() or key in seen:
                continue
            seen.add(key)
            topics.append(t)
        out.append({"session_no": deck["session_no"], "deck_title": deck_title,
                    "topics": topics})
    return out


def taught_digest(course: str, before_session: int, max_per_deck: int = 40) -> str:
    """taught_index() rendered for the prompt — one line per prior session."""
    lines = []
    for entry in taught_index(course, before_session):
        topics = entry["topics"][:max_per_deck]
        more = len(entry["topics"]) - len(topics)
        tail = f" (+{more} more)" if more > 0 else ""
        lines.append(f"  Session {entry['session_no']} — {entry['deck_title']}:\n"
                     f"    {'; '.join(topics)}{tail}")
    return "\n".join(lines)


def taught_titles(course: str, before_session: int) -> list[tuple[int, str]]:
    """(session_no, topic) for every distinct topic already taught IN THIS COURSE — the
    lookup the repetition guardrail compares a new slide's title against."""
    return [(e["session_no"], t)
            for e in taught_index(course, before_session) for t in e["topics"]]


# --------------------------------------------------------------------------- #
# extraction-completeness measure (guideline 2/3: decks must be FULLY extracted)
# --------------------------------------------------------------------------- #
def _source_slide_count(deck: dict) -> int | None:
    """Best-effort ground-truth slide count from the source .pptx, if it is still
    on disk (offline decks in inputs/past_ppts/). Synced decks are extracted from
    in-memory bytes and not kept, so this returns None for them."""
    src = deck.get("source_file")
    if not src:
        return None
    try:
        pat = config.harness()["context"]["past_ppts_glob"]
        base = (config.ROOT / pat).parent
        p = base / src
        if p.exists():
            return len(Presentation(str(p)).slides)
    except Exception:
        pass
    return None


def deck_completeness(deck: dict) -> dict:
    """Deterministic per-deck extraction health from the stored KB JSON.
    A slide with no title AND no body AND no table is treated as 'empty' — a
    likely extraction gap (or a genuinely blank slide)."""
    slides = deck.get("slides", [])
    n = len(slides)
    empty = [s.get("n") for s in slides
             if not (s.get("title") or s.get("body") or s.get("tables"))]
    # Cover (slide 1) and the last two slides are conventionally design/closing
    # slides with little text — an empty one there is NOT an extraction failure.
    # Only INTERIOR empty slides signal genuinely missed content.
    edge = {1, n, n - 1}
    interior_empty = [x for x in empty if x not in edge]
    with_body = sum(1 for s in slides if s.get("body"))
    with_notes = sum(1 for s in slides if s.get("notes"))
    with_tables = sum(1 for s in slides if s.get("tables"))
    src = _source_slide_count(deck)
    dropped = (src - n) if (src is not None) else None
    coverage = round(with_body / n, 3) if n else 0.0

    issues = []
    if n == 0:
        issues.append("no slides extracted")
    if interior_empty:
        issues.append(f"{len(interior_empty)} interior slide(s) with no title/body/table: "
                      f"{interior_empty}")
    if dropped and dropped > 0:
        issues.append(f"extracted {n} of {src} source slides ({dropped} dropped)")

    return {
        "session_no": deck.get("session_no"),
        "source_file": deck.get("source_file"),
        "n_slides": n,
        "source_slides": src,
        "empty_slides": empty,
        "interior_empty_slides": interior_empty,
        "with_body": with_body,
        "with_notes": with_notes,
        "with_tables": with_tables,
        "body_coverage": coverage,
        "ok": not issues,
        "issues": issues,
    }


def completeness_report(course: str | None = None) -> dict:
    """Extraction health across this course's ingested decks.

    `course=None` sweeps EVERY course's store. Diagnostics — the offline eval reports
    extraction health, and its golden belongs to a course whose decks are not in the
    store at all, so it has no single course to ask about.
    """
    if course is None:
        decks = []
        for d in sorted(DECKS_DIR.glob("*/")) if DECKS_DIR.is_dir() else []:
            for f in sorted(d.glob("session_*.json")):
                try:
                    decks.append(json.loads(f.read_text()))
                except Exception:
                    continue
    else:
        decks = load_all_decks(course)
    per = [deck_completeness(d) for d in decks]
    problems = [p for p in per if not p["ok"]]
    return {
        "ok": not problems,
        "decks_checked": len(per),
        "decks_with_issues": len(problems),
        "decks": per,
    }


# --------------------------------------------------------------------------- #
# BM25 lexical RAG retrieval over stored chunks (no API, offline, no deps)
# --------------------------------------------------------------------------- #
_WORD = re.compile(r"[a-z0-9]+")
# Very common words carry no topical signal — dropping them sharpens BM25's IDF.
_STOP = {"the", "a", "an", "and", "or", "of", "to", "in", "is", "are", "for", "on",
         "with", "as", "by", "it", "this", "that", "be", "we", "you", "how", "what",
         "why", "when", "which", "at", "from", "into", "can", "will", "its", "so"}

_K1 = 1.5
_B = 0.75


def _tokens(text: str) -> set[str]:
    return set(_WORD.findall(text.lower()))


def _tok_list(text: str) -> list[str]:
    return [t for t in _WORD.findall(text.lower()) if t not in _STOP]


def slide_detail(slide: dict) -> str:
    """The part of a slide's body that ISN'T just its own title echoed back.

    Section dividers and heading slides extract as a body that repeats the title once or
    twice and says nothing else. They rank fine on a title-word query and then arrive
    carrying no content — which is worse than absent when the block they land in claims
    to show how far a topic was actually taken.
    """
    body = (slide.get("body") or slide.get("notes") or "")
    title = " ".join((slide.get("title") or "").split()).lower()
    if not title:
        return body.strip()
    keep = [ln for ln in body.splitlines()
            if " ".join(ln.split()).lower().strip(":") not in (title, title.strip(":"))]
    return "\n".join(keep).strip()


def rank_slides(decks: list[dict], query: str, top_k: int = 6,
                source: str | None = None, min_detail: int = 0) -> list[dict]:
    """BM25-rank the slides of `decks` against `query`. The ranker, over ANY deck set.

    Split out of retrieve() so the same body-level search can run over a PREREQUISITE's
    decks. It could not before: retrieve() built its corpus from decks_before(course, n)
    and had no way to be pointed anywhere else, so a prerequisite's slide bodies — every
    character of them — were stored and never read. Only their titles ever reached the
    model.

    `source` labels each hit with the course it came from, which matters once hits can
    arrive from more than one place.
    """
    q_terms = [t for t in _tok_list(query)]
    if not q_terms:
        return []

    docs = []  # (session_no, slide, tokens)
    for deck in decks or []:
        for s in deck.get("slides") or []:
            if min_detail:
                # Ask for real content, not a heading that happens to match the query.
                if _BOILERPLATE.match(" ".join((s.get("title") or "").split())):
                    continue
                if len(slide_detail(s)) < min_detail:
                    continue
            blob = " ".join([s.get("title") or "", s.get("body") or "",
                             s.get("notes") or ""])
            toks = _tok_list(blob)
            if toks:
                docs.append((deck.get("session_no"), s, toks))
    if not docs:
        return []

    N = len(docs)
    avgdl = sum(len(t) for _, _, t in docs) / N
    df: dict[str, int] = {}
    for _, _, toks in docs:
        for term in set(toks):
            df[term] = df.get(term, 0) + 1

    q_set = set(q_terms)
    scored = []
    for sn, s, toks in docs:
        dl = len(toks)
        tf: dict[str, int] = {}
        for t in toks:
            if t in q_set:
                tf[t] = tf.get(t, 0) + 1
        if not tf:
            continue
        score = 0.0
        for term, f in tf.items():
            n_q = df.get(term, 0)
            idf = math.log(1 + (N - n_q + 0.5) / (n_q + 0.5))
            score += idf * (f * (_K1 + 1)) / (f + _K1 * (1 - _B + _B * dl / avgdl))
        scored.append((score, sn, s))

    scored.sort(key=lambda x: x[0], reverse=True)
    out = []
    for score, sn, s in scored[:top_k]:
        excerpt = (slide_detail(s) if min_detail
                   else ((s.get("body") or s.get("notes")) or ""))
        hit = {"session_no": sn, "slide": s.get("n"), "title": s.get("title") or "",
               "excerpt": excerpt[:400], "score": round(score, 3)}
        if source:
            hit["source"] = source
        out.append(hit)
    return out


def retrieve(course: str, query: str, session_no: int, top_k: int = 6) -> list[dict]:
    """The most query-relevant slides from THIS course's decks before `session_no`.

    BM25 ranking (Okapi, k1=1.5, b=0.75): rewards rare/distinctive query terms
    (IDF), saturates repeated matches, and normalises by slide length — far
    stronger relevance than raw token overlap. Deterministic, offline, no deps.
    Complements (does not replace) the always-injected per-deck summaries.
    """
    return rank_slides(decks_before(course, session_no), query, top_k=top_k)
